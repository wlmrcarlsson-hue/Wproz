"""Render each slide of a .pptx to HTML from its real geometry.

LibreOffice Impress is not installed in this container, so this reads the
actual generated file with python-pptx and lays every shape out at its true
position and size. Text goes into a box of exactly the shape's dimensions,
so the browser's own wrapping reveals overflow the way PowerPoint would.
Calibri/Cambria are unavailable here and fall back to wider metric cousins,
which makes the overflow check conservative rather than optimistic.
"""
import base64, html, sys
from pptx import Presentation
from pptx.util import Emu

EMU_IN = 914400.0
SCALE = 96  # px per inch
PT = 96 / 72.0  # points -> px

def inches(v):
    return (v or 0) / EMU_IN

def color_of(fmt):
    try:
        if fmt.type is not None and fmt.fore_color.type is not None:
            return "#" + str(fmt.fore_color.rgb)
    except Exception:
        pass
    return None

def run_color(run):
    try:
        if run.font.color and run.font.color.rgb is not None:
            return "#" + str(run.font.color.rgb)
    except Exception:
        pass
    return "#111111"

def render(path, out):
    pr = Presentation(path)
    sw, sh = inches(pr.slide_width), inches(pr.slide_height)
    parts = [f"""<!doctype html><meta charset=utf-8><style>
      body{{margin:0;background:#777;font-family:'Liberation Sans',sans-serif}}
      .slide{{position:relative;width:{sw*SCALE}px;height:{sh*SCALE}px;background:#fff;
             margin:14px auto;overflow:hidden;outline:1px solid #000}}
      .sh{{position:absolute;box-sizing:border-box}}
      .tx{{position:absolute;box-sizing:border-box;overflow:visible}}
      .lbl{{position:absolute;left:0;top:-13px;color:#fff;font:11px monospace}}
      .chart{{position:absolute;box-sizing:border-box;border:2px dashed #888;
              background:#eef1f7;color:#556;font:12px monospace;padding:4px}}
    </style>"""]

    for idx, slide in enumerate(pr.slides, 1):
        bg = "#ffffff"
        try:
            el = slide._element.find(".//{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr", )
            bgel = slide._element.find(".//{http://schemas.openxmlformats.org/presentationml/2006/main}bg")
            if bgel is not None:
                c = bgel.find(".//{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr")
                if c is not None:
                    bg = "#" + c.get("val")
        except Exception:
            pass
        parts.append(f'<div class="slide" data-n="{idx}" style="background:{bg}"><div class="lbl">slide {idx}</div>')
        for sp in slide.shapes:
            x, y = inches(sp.left) * SCALE, inches(sp.top) * SCALE
            w, h = inches(sp.width) * SCALE, inches(sp.height) * SCALE
            rot = sp.rotation or 0
            tf = f"transform:rotate({rot}deg);" if rot else ""

            if sp.shape_type == 13 or sp.__class__.__name__ == "Picture":
                blob = sp.image.blob
                b64 = base64.b64encode(blob).decode()
                parts.append(
                    f'<img class="sh" src="data:{sp.image.content_type};base64,{b64}" '
                    f'style="left:{x}px;top:{y}px;width:{w}px;height:{h}px;{tf}">'
                )
                continue

            if sp.has_chart:
                parts.append(
                    f'<div class="chart" style="left:{x}px;top:{y}px;width:{w}px;height:{h}px">'
                    f'chart</div>'
                )
                continue

            fill = None
            try:
                fill = color_of(sp.fill)
            except Exception:
                pass
            if fill:
                radius = "8px" if "ROUNDED" in str(sp.shape_type) else "0"
                parts.append(
                    f'<div class="sh" style="left:{x}px;top:{y}px;width:{w}px;height:{h}px;'
                    f'background:{fill};border-radius:{radius};{tf}"></div>'
                )

            if not sp.has_text_frame or not sp.text_frame.text.strip():
                continue

            body = []
            for para in sp.text_frame.paragraphs:
                if not para.runs:
                    continue
                r0 = para.runs[0]
                size = (r0.font.size.pt if r0.font.size else 18) * PT
                bold = "font-weight:700;" if r0.font.bold else "font-weight:400;"
                face = r0.font.name or "Calibri"
                fam = "'Liberation Serif',serif" if face in ("Cambria",) else "'Liberation Sans',sans-serif"
                align = {1: "center", 3: "right"}.get(
                    para.alignment.value if para.alignment is not None else 0, "left"
                )
                pPr = para._pPr
                bullet = "• " if (pPr is not None and "buChar" in pPr.xml) else ""
                txt = html.escape("".join(r.text for r in para.runs))
                body.append(
                    f'<div style="font-size:{size:.1f}px;line-height:{size*1.25:.1f}px;font-family:{fam};{bold}'
                    f'color:{run_color(r0)};text-align:{align};margin-bottom:2px">{bullet}{txt}</div>'
                )
            parts.append(
                f'<div class="tx" data-slide="{idx}" style="left:{x}px;top:{y}px;'
                f'width:{w}px;height:{h}px;{tf}">' + "".join(body) + "</div>"
            )
        parts.append("</div>")

    open(out, "w").write("".join(parts))
    print("wrote", out, "slides:", len(pr.slides))

render(sys.argv[1], sys.argv[2])
